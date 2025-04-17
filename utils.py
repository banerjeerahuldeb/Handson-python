import PyPDF2
from docx import Document
from pptx import Presentation
import io

def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_pptx(file):
    ppt = Presentation(file)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(uploaded_file, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(uploaded_file)
    elif file_type == "docx":
        return extract_text_from_docx(uploaded_file)
    elif file_type in ["pptx", "potx"]:
        return extract_text_from_pptx(uploaded_file)
    else:
        return "Unsupported file type."
